import io
import os
from typing import Union

# 所需依赖库（需提前安装）
# pip install python-docx python-pptx openpyxl pdfplumber Pillow

# ========== 模块导入 ==========
try:
    from docx import Document
    from pptx import Presentation
    from openpyxl import load_workbook
    import pdfplumber
except ImportError:
    raise ImportError("请先安装依赖库：python-docx, python-pptx, openpyxl, pdfplumber")

# ========== 异常处理类 ==========
class DocumentParseError(Exception):
    """自定义文档解析异常"""
    def __init__(self, message: str):
        super().__init__(message)
        self.res()

    def res(self)->None:
        print("文档解析异常")


# ========== 核心解析器 ==========
class UniversalDocumentParser:
    def __init__(self):
        self.supported_formats = {
            '.docx': self._parse_word,
            '.pptx': self._parse_ppt,
            '.xlsx': self._parse_excel,
            '.pdf': self._parse_pdf,
            '.doc': self._parse_word,  # 支持Word 97-2003文档
        }

    def parse(self, file_path: str) -> str:
        """
        通用文档解析入口
        :param file_path: 文件路径
        :return: 提取的文本内容
        """
        # 校验文件是否存在
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在：{file_path}")

        # 获取文件扩展名
        ext = os.path.splitext(file_path)[1].lower()

        # 选择解析方法
        parser = self.supported_formats.get(ext)
        if not parser:
            raise DocumentParseError(f"不支持的文件格式：{ext}")

        try:
            return parser(file_path)
        except Exception as e:
            raise DocumentParseError(f"解析失败：{str(e)}") from e

    def _parse_word(self, file_path: str) -> str:
        """解析Word文档"""
        doc = Document(file_path)
        text = []
        
        # 提取段落文本
        for para in doc.paragraphs:
            text.append(para.text)
        
        # 提取表格内容
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                text.append('\t'.join(row_text))
        
        return '\n'.join(text)

    def _parse_ppt(self, file_path: str) -> str:
        """解析PowerPoint文档"""
        prs = Presentation(file_path)
        text = []
        
        # 遍历所有幻灯片
        for slide in prs.slides:
            # 提取形状中的文本
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
                # 处理表格（PPTX中的表格）
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text for cell in row.cells]
                        text.append('\t'.join(row_text))
        
        return '\n'.join(text)

    def _parse_excel(self, file_path: str) -> str:
        """解析Excel文档"""
        wb = load_workbook(file_path)
        text = []
        
        # 遍历所有工作表
        for sheet in wb:
            # 读取最大行和列
            max_row = sheet.max_row
            max_col = sheet.max_column
            
            # 逐行读取数据
            for row in sheet.iter_rows(max_row=max_row, max_col=max_col):
                row_text = [str(cell.value) if cell.value else "" for cell in row]
                text.append('\t'.join(row_text))
        
        return '\n'.join(text)

    def _parse_pdf(self, file_path: str) -> str:
        """解析PDF文档"""
        text = []
        
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                # 提取文本
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
                
                # 提取表格（需要处理表格结构）
                for table in page.extract_tables():
                    for row in table:
                        cleaned_row = [str(cell).replace('\n', ' ') if cell else "" for cell in row]
                        text.append('\t'.join(cleaned_row))
        
        return '\n'.join(text)

# ========== 使用示例 ==========
if __name__ == "__main__":
    parser = UniversalDocumentParser()
    
    # 示例文件路径
    test_files = {
        "Word文档": "example.docx",
        "PPT文档": "ppt.ppt",
        "Excel文档": "data.xlsx",
        "PDF文档": "pdf.pdf"
    }

    for doc_type, path in test_files.items():
        try:
            content = parser.parse(path)
            print(f"========== {doc_type} 解析结果 ==========")
            print(content[:500] + "..." if len(content) > 500 else content)  # 截取前500字符
            print("\n")
        except Exception as e:
            print(f"解析 {doc_type} 失败：{str(e)}")